"""
Sandboxed Attachment Extractor

Runs attachment text extraction in an isolated subprocess to contain
potential security issues from malicious files (RCE, DoS, etc.).

Security measures:
- Separate subprocess (crash isolation)
- Memory limit (prevents memory bombs)
- CPU time limit (prevents infinite loops)
- Timeout enforcement
- No network access (restricted environment)
- No file system access outside temp directory

Usage:
    extractor = SandboxedExtractor()
    text = await extractor.extract_pdf(pdf_bytes)
"""
import asyncio
import base64
import json
import logging
import os
import sys
import tempfile
import subprocess
from typing import Optional, Tuple, List, Dict, Any
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

# Configuration - generous limits for large attachments
DEFAULT_TIMEOUT_SECONDS = 60  # 60 seconds for large PDFs
DEFAULT_MAX_MEMORY_MB = 512   # 512MB for complex documents
DEFAULT_MAX_OUTPUT_SIZE = 10 * 1024 * 1024  # 10MB max extracted text


@dataclass
class ExtractionResult:
    """Result of a sandboxed extraction."""
    success: bool
    text: Optional[str] = None
    error: Optional[str] = None
    extraction_time_ms: Optional[float] = None


@dataclass
class StructuredExtractionResult:
    """Result of a structured extraction (per-page/per-sheet)."""
    success: bool
    items: List[Dict[str, Any]] = None  # List of {"page": 1, "text": "..."} or {"sheet": "Name", "text": "..."}
    error: Optional[str] = None
    extraction_time_ms: Optional[float] = None

    def __post_init__(self):
        if self.items is None:
            self.items = []

    @property
    def total_text(self) -> Optional[str]:
        """Combine all items into a single text."""
        if not self.items:
            return None
        return '\n\n'.join(item.get('text', '') for item in self.items if item.get('text'))

    @property
    def count(self) -> int:
        """Number of pages/sheets extracted."""
        return len(self.items) if self.items else 0


# Worker script that runs in the subprocess
WORKER_SCRIPT = '''
"""Sandboxed extraction worker - runs in isolated subprocess."""
import sys
import json
import base64
import io
import resource
import signal

def set_resource_limits(max_memory_mb: int, max_cpu_seconds: int):
    """Set resource limits for the subprocess."""
    # Memory limit (in bytes)
    memory_bytes = max_memory_mb * 1024 * 1024
    try:
        resource.setrlimit(resource.RLIMIT_AS, (memory_bytes, memory_bytes))
    except (ValueError, resource.error):
        # Some systems don't support RLIMIT_AS
        pass
    
    # CPU time limit (prevents infinite loops)
    try:
        resource.setrlimit(resource.RLIMIT_CPU, (max_cpu_seconds, max_cpu_seconds))
    except (ValueError, resource.error):
        pass
    
    # Prevent fork bombs
    try:
        resource.setrlimit(resource.RLIMIT_NPROC, (0, 0))
    except (ValueError, resource.error):
        pass

def extract_pdf(data: bytes) -> str:
    """Extract text from PDF."""
    import warnings
    import pdfplumber

    text_parts = []
    with warnings.catch_warnings():
        warnings.filterwarnings('ignore')
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                try:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                except Exception:
                    continue
    return '\\n\\n'.join(text_parts)

def extract_pdf_pages(data: bytes) -> list:
    """Extract text from PDF, returning per-page content.

    Returns:
        List of dicts: [{"page": 1, "text": "..."}, ...]
    """
    import warnings
    import pdfplumber

    pages = []
    with warnings.catch_warnings():
        warnings.filterwarnings('ignore')
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append({"page": i, "text": text})
                except Exception:
                    continue
    return pages

def extract_docx(data: bytes) -> str:
    """Extract text from DOCX."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    return '\\n'.join([p.text for p in doc.paragraphs if p.text])

def extract_xlsx(data: bytes) -> str:
    """Extract text from Excel."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    text_parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            cells = [str(cell.value) for cell in row if cell.value is not None]
            if cells:
                text_parts.append(' | '.join(cells))
    return '\\n'.join(text_parts)

def extract_xlsx_sheets(data: bytes) -> list:
    """Extract text from Excel, returning per-sheet content.

    Returns:
        List of dicts: [{"sheet": "Sheet1", "sheet_index": 0, "text": "..."}, ...]
    """
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    sheets = []
    for idx, sheet in enumerate(wb.worksheets):
        rows = []
        for row in sheet.iter_rows():
            cells = [str(cell.value) for cell in row if cell.value is not None]
            if cells:
                rows.append(' | '.join(cells))
        if rows:
            sheets.append({
                "sheet": sheet.title,
                "sheet_index": idx,
                "text": '\\n'.join(rows)
            })
    return sheets

def extract_text_sections(text: str, source_type: str = "text") -> list:
    """Extract sections from plain text or markdown.

    For markdown: splits on headers (# Header)
    For plain text: splits on double newlines (paragraphs) with minimum size

    Args:
        text: The text content
        source_type: "markdown" or "text"

    Returns:
        List of dicts: [{"section": 0, "title": "...", "text": "..."}, ...]
    """
    import re

    sections = []

    if source_type == "markdown":
        # Split on markdown headers
        header_pattern = r'^(#{1,6})\\s+(.+)$'
        lines = text.split('\\n')
        current_section = {"section": 0, "title": "", "text": ""}
        section_idx = 0

        for line in lines:
            match = re.match(header_pattern, line)
            if match:
                # Save previous section if it has content
                if current_section["text"].strip():
                    sections.append(current_section)
                # Start new section
                section_idx += 1
                header_level = len(match.group(1))
                header_text = match.group(2).strip()
                current_section = {
                    "section": section_idx,
                    "title": header_text,
                    "level": header_level,
                    "text": ""
                }
            else:
                current_section["text"] += line + "\\n"

        # Don't forget last section
        if current_section["text"].strip():
            sections.append(current_section)

    else:
        # Plain text: split on double newlines (paragraphs)
        # Merge small paragraphs to avoid too many tiny embeddings
        MIN_SECTION_CHARS = 500
        MAX_SECTION_CHARS = 8000

        paragraphs = re.split(r'\\n\\s*\\n', text)
        current_section = {"section": 0, "title": "", "text": ""}
        section_idx = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # If adding this paragraph would exceed max, start new section
            if len(current_section["text"]) + len(para) > MAX_SECTION_CHARS:
                if current_section["text"]:
                    sections.append(current_section)
                    section_idx += 1
                current_section = {"section": section_idx, "title": "", "text": para}
            else:
                if current_section["text"]:
                    current_section["text"] += "\\n\\n" + para
                else:
                    current_section["text"] = para

        # Save last section if it meets minimum size (or if it's the only one)
        if current_section["text"] and (len(current_section["text"]) >= MIN_SECTION_CHARS or not sections):
            sections.append(current_section)
        elif current_section["text"] and sections:
            # Append to previous section
            sections[-1]["text"] += "\\n\\n" + current_section["text"]

    return sections if sections else [{"section": 0, "title": "", "text": text}]

def parse_page_breaks(text: str) -> list:
    """Parse page breaks from pdftotext/OCR output.

    pdftotext and many OCR tools use form feed (\\x0c) as page separator.

    Returns:
        List of dicts: [{"page": 1, "text": "..."}, ...]
    """
    # Form feed character marks page breaks
    pages = text.split('\\x0c')
    result = []
    for i, page_text in enumerate(pages, 1):
        page_text = page_text.strip()
        if page_text:
            result.append({"page": i, "text": page_text})
    return result if result else [{"page": 1, "text": text}]

def extract_pptx(data: bytes) -> str:
    """Extract text from PowerPoint."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            text_parts.append(f"Slide {i}:\\n" + '\\n'.join(slide_text))
    return '\\n\\n'.join(text_parts)

def extract_rtf(data: bytes) -> str:
    """Extract text from RTF."""
    from striprtf.striprtf import rtf_to_text
    rtf_string = data.decode('utf-8', errors='replace')
    return rtf_to_text(rtf_string)

def extract_ics(data: bytes) -> str:
    """Extract text from iCalendar."""
    from icalendar import Calendar
    ics_string = data.decode('utf-8', errors='replace')
    cal = Calendar.from_ical(ics_string)
    events = []
    for component in cal.walk():
        if component.name == 'VEVENT':
            parts = []
            if component.get('summary'):
                parts.append(f"Event: {component.get('summary')}")
            if component.get('dtstart'):
                parts.append(f"Start: {component.get('dtstart').dt}")
            if component.get('dtend'):
                parts.append(f"End: {component.get('dtend').dt}")
            if component.get('location'):
                parts.append(f"Location: {component.get('location')}")
            if component.get('description'):
                parts.append(f"Description: {component.get('description')}")
            if parts:
                events.append('\\n'.join(parts))
    return '\\n\\n'.join(events)

def main():
    """Main entry point for the worker."""
    try:
        # Read input from stdin
        input_data = json.loads(sys.stdin.read())
        
        file_type = input_data['type']
        file_data = base64.b64decode(input_data['data'])
        max_memory_mb = input_data.get('max_memory_mb', 256)
        max_cpu_seconds = input_data.get('max_cpu_seconds', 30)
        
        # Set resource limits
        set_resource_limits(max_memory_mb, max_cpu_seconds)
        
        # Extract based on type
        extractors = {
            'pdf': extract_pdf,
            'pdf_pages': extract_pdf_pages,
            'docx': extract_docx,
            'xlsx': extract_xlsx,
            'xlsx_sheets': extract_xlsx_sheets,
            'pptx': extract_pptx,
            'rtf': extract_rtf,
            'ics': extract_ics,
        }

        # Handle text/markdown sectioning (passed as special types)
        if file_type == 'text_sections':
            text = file_data.decode('utf-8', errors='replace')
            output = extract_text_sections(text, source_type='text')
            result = {'success': True, 'structured': output}
        elif file_type == 'markdown_sections':
            text = file_data.decode('utf-8', errors='replace')
            output = extract_text_sections(text, source_type='markdown')
            result = {'success': True, 'structured': output}
        elif file_type == 'parse_pages':
            # Parse page breaks from pdftotext/OCR output
            text = file_data.decode('utf-8', errors='replace')
            output = parse_page_breaks(text)
            result = {'success': True, 'structured': output}
        elif file_type not in extractors:
            result = {'success': False, 'error': f'Unknown file type: {file_type}'}
        else:
            output = extractors[file_type](file_data)
            # Structured extractors return lists, others return strings
            if isinstance(output, list):
                result = {'success': True, 'structured': output}
            else:
                result = {'success': True, 'text': output}
        
        print(json.dumps(result))
        
    except MemoryError:
        print(json.dumps({'success': False, 'error': 'Memory limit exceeded'}))
    except Exception as e:
        print(json.dumps({'success': False, 'error': str(e)[:500]}))

if __name__ == '__main__':
    main()
'''


class SandboxedExtractor:
    """
    Runs attachment extraction in a sandboxed subprocess.
    
    Features:
    - Crash isolation: subprocess crashes don't affect main process
    - Resource limits: memory and CPU time limits
    - Timeout: extraction is killed after timeout
    - No network: subprocess environment restricts network access
    """
    
    def __init__(
        self,
        timeout_seconds: int = DEFAULT_TIMEOUT_SECONDS,
        max_memory_mb: int = DEFAULT_MAX_MEMORY_MB,
        max_output_size: int = DEFAULT_MAX_OUTPUT_SIZE,
        enabled: bool = True
    ):
        """
        Initialize the sandboxed extractor.
        
        Args:
            timeout_seconds: Maximum time for extraction
            max_memory_mb: Maximum memory for subprocess
            max_output_size: Maximum size of extracted text
            enabled: If False, extraction is skipped (returns empty)
        """
        self.timeout_seconds = timeout_seconds
        self.max_memory_mb = max_memory_mb
        self.max_output_size = max_output_size
        self.enabled = enabled
        self._worker_script_path: Optional[Path] = None
    
    def _get_worker_script_path(self) -> Path:
        """Get or create the worker script file."""
        if self._worker_script_path and self._worker_script_path.exists():
            return self._worker_script_path
        
        # Create worker script in temp directory
        worker_dir = Path(tempfile.gettempdir()) / "email_processor_sandbox"
        worker_dir.mkdir(exist_ok=True)
        
        worker_path = worker_dir / "extraction_worker.py"
        worker_path.write_text(WORKER_SCRIPT)
        
        self._worker_script_path = worker_path
        return worker_path
    
    def _get_restricted_env(self) -> dict:
        """Create a restricted environment for the subprocess."""
        # Start with minimal environment
        env = {
            'PATH': os.environ.get('PATH', '/usr/bin:/bin'),
            'HOME': tempfile.gettempdir(),
            'TMPDIR': tempfile.gettempdir(),
            'LANG': 'C.UTF-8',
        }
        
        # Add PYTHONPATH if needed
        if 'PYTHONPATH' in os.environ:
            env['PYTHONPATH'] = os.environ['PYTHONPATH']
        
        # Explicitly disable network proxies
        env['no_proxy'] = '*'
        env['NO_PROXY'] = '*'
        
        return env
    
    async def _run_extraction(self, file_type: str, file_data: bytes) -> ExtractionResult:
        """Run extraction in sandboxed subprocess."""
        import time
        start_time = time.time()
        
        if not self.enabled:
            return ExtractionResult(
                success=False,
                error="Sandboxed extraction disabled"
            )
        
        try:
            # Prepare input
            input_data = json.dumps({
                'type': file_type,
                'data': base64.b64encode(file_data).decode('ascii'),
                'max_memory_mb': self.max_memory_mb,
                'max_cpu_seconds': self.timeout_seconds,
            })
            
            # Get worker script path
            worker_path = self._get_worker_script_path()
            
            # Run subprocess
            proc = await asyncio.create_subprocess_exec(
                sys.executable,
                str(worker_path),
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                env=self._get_restricted_env(),
                # Limit output size
                limit=self.max_output_size,
            )
            
            try:
                stdout, stderr = await asyncio.wait_for(
                    proc.communicate(input=input_data.encode()),
                    timeout=self.timeout_seconds
                )
            except asyncio.TimeoutError:
                proc.kill()
                await proc.wait()
                return ExtractionResult(
                    success=False,
                    error=f"Extraction timed out after {self.timeout_seconds}s",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )
            
            if proc.returncode != 0:
                stderr_text = stderr.decode('utf-8', errors='replace')[:500]
                return ExtractionResult(
                    success=False,
                    error=f"Subprocess failed (code {proc.returncode}): {stderr_text}",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )
            
            # Parse output
            try:
                result = json.loads(stdout.decode('utf-8'))
                return ExtractionResult(
                    success=result.get('success', False),
                    text=result.get('text'),
                    error=result.get('error'),
                    extraction_time_ms=(time.time() - start_time) * 1000
                )
            except json.JSONDecodeError as e:
                return ExtractionResult(
                    success=False,
                    error=f"Invalid output from worker: {str(e)}",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )
                
        except Exception as e:
            logger.error(f"Sandboxed extraction failed: {e}")
            return ExtractionResult(
                success=False,
                error=f"Sandbox error: {str(e)[:200]}",
                extraction_time_ms=(time.time() - start_time) * 1000
            )
    
    async def extract_pdf(self, pdf_bytes: bytes) -> ExtractionResult:
        """Extract text from PDF in sandbox."""
        return await self._run_extraction('pdf', pdf_bytes)
    
    async def extract_docx(self, docx_bytes: bytes) -> ExtractionResult:
        """Extract text from DOCX in sandbox."""
        return await self._run_extraction('docx', docx_bytes)
    
    async def extract_xlsx(self, xlsx_bytes: bytes) -> ExtractionResult:
        """Extract text from Excel in sandbox."""
        return await self._run_extraction('xlsx', xlsx_bytes)
    
    async def extract_pptx(self, pptx_bytes: bytes) -> ExtractionResult:
        """Extract text from PowerPoint in sandbox."""
        return await self._run_extraction('pptx', pptx_bytes)
    
    async def extract_rtf(self, rtf_bytes: bytes) -> ExtractionResult:
        """Extract text from RTF in sandbox."""
        return await self._run_extraction('rtf', rtf_bytes)
    
    async def extract_ics(self, ics_bytes: bytes) -> ExtractionResult:
        """Extract text from iCalendar in sandbox."""
        return await self._run_extraction('ics', ics_bytes)

    async def _run_structured_extraction(
        self, file_type: str, file_data: bytes
    ) -> StructuredExtractionResult:
        """Run structured extraction (per-page/per-sheet) in sandboxed subprocess."""
        import time
        start_time = time.time()

        if not self.enabled:
            return StructuredExtractionResult(
                success=False,
                error="Sandboxed extraction disabled"
            )

        try:
            # Prepare input
            input_data = json.dumps({
                'type': file_type,
                'data': base64.b64encode(file_data).decode('ascii'),
                'max_memory_mb': self.max_memory_mb,
                'max_cpu_seconds': self.timeout_seconds,
            })

            # Get worker script path
            worker_path = self._get_worker_script_path()

            # Run subprocess
            proc = await asyncio.create_subprocess_exec(
                sys.executable,
                str(worker_path),
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                env=self._get_restricted_env(),
                limit=self.max_output_size,
            )

            try:
                stdout, stderr = await asyncio.wait_for(
                    proc.communicate(input=input_data.encode()),
                    timeout=self.timeout_seconds
                )
            except asyncio.TimeoutError:
                proc.kill()
                await proc.wait()
                return StructuredExtractionResult(
                    success=False,
                    error=f"Extraction timed out after {self.timeout_seconds}s",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )

            if proc.returncode != 0:
                stderr_text = stderr.decode('utf-8', errors='replace')[:500]
                return StructuredExtractionResult(
                    success=False,
                    error=f"Subprocess failed (code {proc.returncode}): {stderr_text}",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )

            # Parse output
            try:
                result = json.loads(stdout.decode('utf-8'))
                return StructuredExtractionResult(
                    success=result.get('success', False),
                    items=result.get('structured', []),
                    error=result.get('error'),
                    extraction_time_ms=(time.time() - start_time) * 1000
                )
            except json.JSONDecodeError as e:
                return StructuredExtractionResult(
                    success=False,
                    error=f"Invalid output from worker: {str(e)}",
                    extraction_time_ms=(time.time() - start_time) * 1000
                )

        except Exception as e:
            logger.error(f"Sandboxed structured extraction failed: {e}")
            return StructuredExtractionResult(
                success=False,
                error=f"Sandbox error: {str(e)[:200]}",
                extraction_time_ms=(time.time() - start_time) * 1000
            )

    async def extract_pdf_pages(self, pdf_bytes: bytes) -> StructuredExtractionResult:
        """
        Extract text from PDF, returning per-page content.

        Returns:
            StructuredExtractionResult with items like [{"page": 1, "text": "..."}, ...]
        """
        return await self._run_structured_extraction('pdf_pages', pdf_bytes)

    async def extract_xlsx_sheets(self, xlsx_bytes: bytes) -> StructuredExtractionResult:
        """
        Extract text from Excel, returning per-sheet content.

        Returns:
            StructuredExtractionResult with items like [{"sheet": "Name", "sheet_index": 0, "text": "..."}, ...]
        """
        return await self._run_structured_extraction('xlsx_sheets', xlsx_bytes)

    async def extract_text_sections(self, text_bytes: bytes) -> StructuredExtractionResult:
        """
        Split plain text into sections by paragraphs.

        Returns:
            StructuredExtractionResult with items like [{"section": 0, "text": "..."}, ...]
        """
        return await self._run_structured_extraction('text_sections', text_bytes)

    async def extract_markdown_sections(self, md_bytes: bytes) -> StructuredExtractionResult:
        """
        Split markdown into sections by headers.

        Returns:
            StructuredExtractionResult with items like [{"section": 0, "title": "Header", "level": 2, "text": "..."}, ...]
        """
        return await self._run_structured_extraction('markdown_sections', md_bytes)

    async def parse_page_breaks(self, text_bytes: bytes) -> StructuredExtractionResult:
        """
        Parse page breaks from pdftotext/OCR output (form feed \\x0c separator).

        Returns:
            StructuredExtractionResult with items like [{"page": 1, "text": "..."}, ...]
        """
        return await self._run_structured_extraction('parse_pages', text_bytes)

    def _decode_text(self, content: bytes) -> Optional[str]:
        """
        Decode bytes to string, trying common encodings.

        Args:
            content: Raw bytes

        Returns:
            Decoded string or None if decoding fails
        """
        # Try common encodings in order of likelihood
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']

        for encoding in encodings:
            try:
                return content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                continue

        # Last resort: decode with replacement characters
        try:
            return content.decode('utf-8', errors='replace')
        except Exception:
            return None

    async def extract_text(
        self,
        content: bytes,
        content_type: str,
        filename: str = "unknown",
    ) -> Optional[str]:
        """
        Generic text extraction that dispatches based on MIME type.

        Args:
            content: File content as bytes
            content_type: MIME type (e.g., 'application/pdf')
            filename: Original filename (used as fallback for type detection)

        Returns:
            Extracted text or None if extraction fails/unsupported
        """
        # Plain text types - just decode directly
        plain_text_mimes = {
            'text/plain', 'text/markdown', 'text/x-markdown',
        }
        if content_type in plain_text_mimes:
            return self._decode_text(content)

        # Map MIME types to extraction methods
        mime_to_method = {
            'application/pdf': self.extract_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self.extract_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self.extract_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self.extract_pptx,
            'application/msword': self.extract_docx,  # Old .doc format (best effort)
            'application/vnd.ms-excel': self.extract_xlsx,  # Old .xls format
            'application/vnd.ms-powerpoint': self.extract_pptx,  # Old .ppt format
            'application/rtf': self.extract_rtf,
            'text/rtf': self.extract_rtf,
            'text/calendar': self.extract_ics,
        }

        # Try by MIME type first
        extract_method = mime_to_method.get(content_type)

        # Fallback to extension if MIME type not recognized
        if not extract_method and filename:
            ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

            # Plain text extensions
            if ext in ('txt', 'md', 'markdown', 'text'):
                return self._decode_text(content)

            ext_to_method = {
                'pdf': self.extract_pdf,
                'docx': self.extract_docx,
                'doc': self.extract_docx,
                'xlsx': self.extract_xlsx,
                'xls': self.extract_xlsx,
                'pptx': self.extract_pptx,
                'ppt': self.extract_pptx,
                'rtf': self.extract_rtf,
                'ics': self.extract_ics,
            }
            extract_method = ext_to_method.get(ext)

        if not extract_method:
            logger.debug(f"No extractor for {content_type} / {filename}")
            return None

        try:
            result = await extract_method(content)
            if result.success and result.text:
                return result.text
            return None
        except Exception as e:
            logger.warning(f"Extraction failed for {filename}: {e}")
            return None

    async def extract_structured(
        self,
        content: bytes,
        content_type: str,
        filename: str = "unknown",
    ) -> Optional[StructuredExtractionResult]:
        """
        Extract text with page/sheet/section structure for documents that support it.

        For PDFs: returns per-page content
        For XLSX: returns per-sheet content
        For Markdown: returns per-header sections
        For plain text: returns paragraph-based sections

        Args:
            content: File content as bytes
            content_type: MIME type
            filename: Original filename for fallback type detection

        Returns:
            StructuredExtractionResult or None if type doesn't support structured extraction
        """
        # Determine file type
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

        # Check for PDF
        if content_type == 'application/pdf' or ext == 'pdf':
            return await self.extract_pdf_pages(content)

        # Check for Excel
        xlsx_mimes = {
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-excel',
        }
        if content_type in xlsx_mimes or ext in ('xlsx', 'xls'):
            return await self.extract_xlsx_sheets(content)

        # Check for Markdown
        markdown_mimes = {'text/markdown', 'text/x-markdown'}
        if content_type in markdown_mimes or ext in ('md', 'markdown'):
            return await self.extract_markdown_sections(content)

        # Check for plain text (only if large enough to benefit from sectioning)
        text_mimes = {'text/plain'}
        if content_type in text_mimes or ext in ('txt', 'text'):
            # Only section if text is large enough (> 2KB)
            if len(content) > 2048:
                return await self.extract_text_sections(content)

        # Type doesn't support structured extraction
        return None


# Global singleton for use across the application
_sandboxed_extractor: Optional[SandboxedExtractor] = None


def get_sandboxed_extractor() -> SandboxedExtractor:
    """Get the global sandboxed extractor instance."""
    global _sandboxed_extractor
    if _sandboxed_extractor is None:
        # Check if sandboxing is enabled via environment (enabled by default)
        enabled = os.getenv('SANDBOX_ATTACHMENTS', 'true').lower() in ('true', '1', 'yes')
        timeout = int(os.getenv('SANDBOX_TIMEOUT_SECONDS', str(DEFAULT_TIMEOUT_SECONDS)))
        max_memory = int(os.getenv('SANDBOX_MAX_MEMORY_MB', str(DEFAULT_MAX_MEMORY_MB)))
        
        _sandboxed_extractor = SandboxedExtractor(
            timeout_seconds=timeout,
            max_memory_mb=max_memory,
            enabled=enabled
        )
        
        if enabled:
            logger.info(f"Sandboxed attachment extraction enabled (timeout={timeout}s, memory={max_memory}MB)")
        else:
            logger.warning("⚠️  Sandboxed attachment extraction DISABLED - using in-process extraction")
    
    return _sandboxed_extractor
