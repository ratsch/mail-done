"""
Email processing following Inbox Zero's approach.
Handles parsing, HTML to Markdown conversion, and attachment text extraction.

Security: Attachment extraction runs in a sandboxed subprocess to contain
potential RCE/DoS attacks from malicious files. See sandboxed_extractor.py.
"""
import email
from email import policy
from email.utils import parsedate_to_datetime
from email.message import Message
from datetime import datetime, timezone
from markdownify import markdownify
from bs4 import BeautifulSoup
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from icalendar import Calendar
import io
import os
from typing import List, Optional, Tuple
import logging
import re

from .models import ProcessedEmail, AttachmentInfo
from .attachment_detection import is_attachment_part
from .sandboxed_extractor import get_sandboxed_extractor, ExtractionResult

logger = logging.getLogger(__name__)

# Check if sandboxing is enabled (default: True for security)
# WARNING: When disabled, malicious attachments could execute arbitrary code!
USE_SANDBOXED_EXTRACTION = os.getenv('SANDBOX_ATTACHMENTS', 'true').lower() in ('true', '1', 'yes')

if not USE_SANDBOXED_EXTRACTION:
    logger.error("🚨 SECURITY WARNING: SANDBOX_ATTACHMENTS=false - attachment extraction runs IN-PROCESS!")
    logger.error("🚨 This is a security risk. Malicious attachments could execute arbitrary code.")


def _log_sandbox_failure(file_type: str, error: str, size_bytes: int = 0):
    """Log sandbox extraction failure with high visibility."""
    size_str = f" ({size_bytes / 1024:.1f}KB)" if size_bytes else ""
    logger.error(f"🔴 SANDBOX EXTRACTION FAILED [{file_type}]{size_str}: {error}")
    logger.error(f"🔴 Attachment text will NOT be extracted (no fallback to in-process for security)")
    # Also log at warning level for monitoring tools that filter by level
    logger.warning(f"Sandbox extraction failed for {file_type}: {error}")


class EmailProcessor:
    """Mirrors inbox-zero's email processing logic"""
    
    def __init__(self):
        # Use same markdown conversion settings as inbox-zero
        self.markdown_options = {
            'heading_style': 'ATX',
            'bullets': '-',
            'strong_em_symbol': '**',
            'strip': ['script', 'style'],  # Security: remove scripts
        }
    
    async def process(self, raw_email: bytes, uid: str, folder: str = "INBOX", return_headers: bool = False) -> ProcessedEmail:
        """
        Process email like inbox-zero does.
        
        Args:
            raw_email: Raw RFC822 email bytes
            uid: IMAP UID of the email
            folder: IMAP folder name (e.g., "INBOX", "Sent", "Archive")
            
        Returns:
            ProcessedEmail with normalized content
        """
        try:
            # Use modern email policy for better handling of:
            # - Encoded filenames (=?utf-8?Q?...?=)
            # - Unicode content
            # - Consistent parameter parsing
            msg = email.message_from_bytes(raw_email, policy=policy.default)
        except Exception as e:
            logger.error(f"Failed to parse email {uid}: {e}")
            # Return minimal email with error info
            return ProcessedEmail(
                uid=uid,
                message_id=f"<error-{uid}>",
                thread_id=None,
                subject=f"[PARSE ERROR] Email {uid}",
                from_address="unknown@error",
                from_name=None,
                to_addresses=[],
                to_names=[],
                cc_addresses=[],
                date=datetime.now(),
                body_markdown=f"Failed to parse email: {str(e)[:200]}",
                body_text=f"Failed to parse email: {str(e)[:200]}",
                sender_domain="error",
            )
        
        # Extract ALL headers (for preprocessing)
        # Convert Header objects to strings (some headers may be Header objects)
        raw_headers = {}
        for key, value in msg.items():
            if isinstance(value, str):
                raw_headers[key] = value
            else:
                # Convert email.header.Header objects to string
                raw_headers[key] = str(value)
        
        # Extract basic headers with validation
        subject = self._decode_header(msg.get('Subject', ''))
        if not subject or subject.strip() == '':
            subject = "(No Subject)"
        
        from_address, from_name = self._extract_email_with_name(msg.get('From', ''))
        if not from_address:
            logger.warning(f"Email {uid} has no From address")
            from_address = "unknown@unknown"
        
        to_addresses, to_names = self._extract_emails_with_names(msg.get_all('To', []))
        if not to_addresses:
            logger.warning(f"Email {uid} has no To addresses")
            to_addresses = ["unknown@unknown"]
        
        # Extract CC addresses
        cc_addresses, _ = self._extract_emails_with_names(msg.get_all('Cc', []))
        
        # Parse date with validation
        date = self._parse_date_safe(msg.get('Date'), uid)
        
        # Extract body (follow inbox-zero's preference order)
        body_html, body_text = self._extract_body(msg, uid)
        
        # Convert to markdown (inbox-zero pattern) with validation
        body_markdown = self._html_to_markdown(body_html) if body_html else (body_text or "")
        if not body_markdown or body_markdown.strip() == '':
            body_markdown = "(Empty email body)"
        
        # Store body_text for rule matching and other uses
        stored_body_text = body_text or ""
        
        # Extract attachments (inbox-zero feature)
        attachment_texts, attachment_infos = await self._extract_attachments(msg)
        
        # Extract sender domain
        sender_domain = from_address.split('@')[-1] if '@' in from_address else ""
        
        # Extract threading headers (for reply tracking in Phase 2)
        # Use References header or In-Reply-To, fallback to Message-ID
        references = msg.get('References')
        thread_id = (
            references.split()[0] if references
            else msg.get('In-Reply-To') 
            or msg.get('Message-ID')
        )
        
        return ProcessedEmail(
            uid=uid,
            message_id=msg.get('Message-ID'),
            thread_id=thread_id,
            references=references,  # Include full References header
            subject=subject,
            from_address=from_address,
            from_name=from_name,
            to_addresses=to_addresses,
            to_names=to_names,
            cc_addresses=cc_addresses,
            date=date,
            body_markdown=body_markdown,
            body_text=stored_body_text,
            attachment_texts=attachment_texts,
            attachment_info=attachment_infos,  # Include attachment metadata
            sender_domain=sender_domain,
            has_attachments=len(attachment_infos) > 0,
            attachment_count=len(attachment_infos),
            raw_headers=raw_headers,  # Include for preprocessing
            folder=folder,
            flags=[],
        )
    
    def _decode_header(self, header: str) -> str:
        """Decode email header (handles encoding with fallbacks for unknown charsets)"""
        if not header:
            return ""
        
        from email.header import decode_header
        decoded_parts = []
        
        for part, encoding in decode_header(header):
            if isinstance(part, bytes):
                # Try specified encoding first, with fallbacks for unknown charsets
                if encoding:
                    # Map common unknown/non-standard encodings to known ones
                    encoding_map = {
                        'x-euc-jp': 'euc-jp',
                        'x-sjis': 'shift-jis',
                        'x-gb2312': 'gb2312',
                        'x-big5': 'big5'
                    }
                    encoding = encoding_map.get(encoding.lower(), encoding)
                    
                    # Detect MIME types incorrectly marked as encodings
                    # (e.g., "text/html", "text/plain" are MIME types, not charsets)
                    if encoding.lower().startswith('text/') or encoding.lower().startswith('application/'):
                        # This is a MIME type, not an encoding - use utf-8
                        logger.debug(f"Detected MIME type '{encoding}' as encoding, using utf-8")
                        encoding = 'utf-8'
                    
                    try:
                        decoded_parts.append(part.decode(encoding, errors='replace'))
                    except (LookupError, UnicodeDecodeError):
                        # Unknown encoding - try common fallbacks
                        for fallback in ['utf-8', 'latin-1', 'ascii']:
                            try:
                                decoded_parts.append(part.decode(fallback, errors='replace'))
                                break
                            except:
                                continue
                        else:
                            # All failed, force to string
                            decoded_parts.append(str(part, errors='replace'))
                else:
                    # No encoding specified, try utf-8
                    try:
                        decoded_parts.append(part.decode('utf-8', errors='replace'))
                    except:
                        decoded_parts.append(str(part, errors='replace'))
            else:
                decoded_parts.append(str(part))
        
        return ''.join(decoded_parts)
    
    def _parse_date_safe(self, date_str: Optional[str], uid: str) -> datetime:
        """
        Parse email date with validation and fallbacks.
        Handles invalid dates, future dates, malformed RFC 2822, and invalid timezone offsets.
        """
        if not date_str:
            logger.debug(f"Email {uid}: No date header, using current time")
            return datetime.now(timezone.utc)
        
        try:
            parsed_date = parsedate_to_datetime(date_str)
            
            # Validate and normalize timezone offset
            # PostgreSQL requires timezone offsets between -12:00 and +14:00
            if parsed_date.tzinfo is not None:
                offset = parsed_date.utcoffset()
                if offset is not None:
                    total_seconds = offset.total_seconds()
                    # Check if offset is out of valid range (-12h to +14h)
                    if total_seconds < -43200 or total_seconds > 50400:  # -12h to +14h in seconds
                        logger.warning(f"Email {uid}: Invalid timezone offset {offset} (valid range: -12:00 to +14:00), normalizing to UTC")
                        # Convert to UTC by removing timezone info and adding UTC
                        parsed_date = parsed_date.replace(tzinfo=None)
                        parsed_date = parsed_date.replace(tzinfo=timezone.utc)
                else:
                    # No offset but has tzinfo, assume UTC
                    parsed_date = parsed_date.replace(tzinfo=timezone.utc)
            else:
                # No timezone info, assume UTC
                parsed_date = parsed_date.replace(tzinfo=timezone.utc)
            
            # Validate: reject dates too far in future (> 1 day)
            now = datetime.now(timezone.utc)
            if parsed_date > now:
                time_diff = (parsed_date - now).total_seconds()
                if time_diff > 86400:  # More than 1 day in future
                    logger.warning(f"Email {uid}: Date {parsed_date} is {time_diff/3600:.1f}h in future, using current time")
                    return datetime.now(timezone.utc)
            
            # Validate: reject dates too far in past (before email was invented, ~1970)
            if parsed_date.year < 1970:
                logger.warning(f"Email {uid}: Date {parsed_date} is before 1970, using current time")
                return datetime.now(timezone.utc)
            
            return parsed_date
            
        except Exception as e:
            logger.warning(f"Email {uid}: Failed to parse date '{date_str}': {e}, using current time")
            return datetime.now(timezone.utc)
    
    def _decode_body_safe(self, payload: bytes, charset: str, uid: str) -> str:
        """
        Decode email body with fallbacks for unknown/invalid encodings.
        Handles encoding errors gracefully.
        """
        if not charset:
            charset = 'utf-8'
        
        # Map common unknown/non-standard encodings to known ones
        encoding_map = {
            'x-unknown': 'utf-8',
            'ansi_x3.110-1983': 'latin-1',  # Tamil encoding, fallback to latin-1
            'x-euc-jp': 'euc-jp',
            'x-sjis': 'shift-jis',
            'x-gb2312': 'gb2312',
            'x-big5': 'big5'
        }
        
        charset_lower = charset.lower()
        if charset_lower in encoding_map:
            logger.debug(f"Email {uid}: Mapping unknown encoding '{charset}' to '{encoding_map[charset_lower]}'")
            charset = encoding_map[charset_lower]
        
        # Try the specified encoding first
        try:
            return payload.decode(charset, errors='replace')
        except (LookupError, UnicodeDecodeError) as e:
            # Unknown encoding or decode error - try fallbacks
            logger.debug(f"Email {uid}: Failed to decode with '{charset}': {e}, trying fallbacks")
            for fallback in ['utf-8', 'latin-1', 'cp1252', 'ascii']:
                try:
                    return payload.decode(fallback, errors='replace')
                except:
                    continue
            
            # All failed, force decode with utf-8 and replace errors
            logger.warning(f"Email {uid}: All encoding fallbacks failed for charset '{charset}', forcing utf-8 with error replacement")
            return payload.decode('utf-8', errors='replace')
    
    def _extract_email_address(self, header: str) -> str:
        """Extract email address from header (backward compatibility)"""
        if not header:
            return ""
        
        from email.utils import parseaddr
        name, addr = parseaddr(header)
        return addr
    
    def _extract_email_with_name(self, header: str) -> tuple[str, Optional[str]]:
        """
        Extract email address and display name separately.
        
        Args:
            header: Email header (e.g., "John Doe <john@example.com>")
            
        Returns:
            (email_address, display_name) tuple
        """
        if not header:
            return ("", None)
        
        from email.utils import parseaddr
        name, addr = parseaddr(self._decode_header(header))
        
        # Clean up name (remove extra quotes, whitespace)
        if name:
            name = name.strip('"').strip("'").strip()
        
        return (addr, name if name else None)
    
    def _extract_email_addresses(self, headers: List[str]) -> List[str]:
        """Extract multiple email addresses (backward compatibility)"""
        addresses = []
        for header in headers:
            if header:
                from email.utils import getaddresses
                for name, addr in getaddresses([header]):
                    if addr:
                        addresses.append(addr)
        return addresses
    
    def _extract_emails_with_names(self, headers: List[str]) -> tuple[List[str], List[str]]:
        """
        Extract multiple email addresses with display names.
        
        Args:
            headers: List of email headers
            
        Returns:
            (email_addresses, display_names) tuple
        """
        addresses = []
        names = []
        
        for header in headers:
            if header:
                from email.utils import getaddresses
                for name, addr in getaddresses([self._decode_header(header)]):
                    if addr:
                        addresses.append(addr)
                        # Clean up name
                        clean_name = name.strip('"').strip("'").strip() if name else None
                        names.append(clean_name if clean_name else "")
        
        return (addresses, names)
    
    def _extract_body(self, msg: Message, uid: str = "unknown") -> Tuple[Optional[str], Optional[str]]:
        """
        Extract HTML and text body from email.
        
        Args:
            msg: Email message object
            uid: Email UID for logging (optional, defaults to "unknown")
            
        Returns: (html_body, text_body)
        """
        html_body = None
        text_body = None
        
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                disposition = str(part.get('Content-Disposition', ''))
                
                # Skip attachments
                if 'attachment' in disposition:
                    continue
                
                if content_type == 'text/plain' and text_body is None:
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        text_body = self._decode_body_safe(payload, charset, uid)
                
                elif content_type == 'text/html' and html_body is None:
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        html_body = self._decode_body_safe(payload, charset, uid)
        else:
            # Single part message
            content_type = msg.get_content_type()
            payload = msg.get_payload(decode=True)
            
            if payload:
                charset = msg.get_content_charset() or 'utf-8'
                decoded = self._decode_body_safe(payload, charset, uid)
                
                if content_type == 'text/html':
                    html_body = decoded
                else:
                    text_body = decoded
        
        return html_body, text_body
    
    def _html_to_markdown(self, html: str) -> str:
        """Convert HTML to markdown like inbox-zero"""
        if not html:
            return ""
        
        # Clean HTML first (inbox-zero approach)
        soup = BeautifulSoup(html, 'html.parser')
        
        # Remove scripts and styles (security from inbox-zero)
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Remove tracking pixels and other unwanted elements
        for tag in soup.find_all(['img'], {'width': '1', 'height': '1'}):
            tag.decompose()
        
        # Convert to markdown
        try:
            markdown = markdownify(str(soup), **self.markdown_options)
            # Clean up excessive newlines
            while '\n\n\n' in markdown:
                markdown = markdown.replace('\n\n\n', '\n\n')
            return markdown.strip()
        except Exception as e:
            logger.error(f"Failed to convert HTML to markdown: {e}")
            # Fallback to text extraction
            return soup.get_text()
    
    async def _extract_attachments(self, msg: Message, depth: int = 0) -> Tuple[List[str], List[AttachmentInfo]]:
        """
        Extract text from attachments (inbox-zero feature).
        Includes error handling per attachment and size limits.
        
        Args:
            msg: Email message to process
            depth: Recursion depth (for nested message/rfc822 parts)
        """
        texts = []
        infos = []
        
        # Size limit for attachment processing (50MB)
        MAX_ATTACHMENT_SIZE = 50 * 1024 * 1024
        
        # Limit recursion depth to prevent infinite loops
        MAX_RECURSION_DEPTH = 5
        if depth > MAX_RECURSION_DEPTH:
            logger.warning(f"Max recursion depth ({MAX_RECURSION_DEPTH}) reached in attachment extraction")
            return texts, infos
        
        for part in msg.walk():
            content_type = part.get_content_type()
            filename = part.get_filename()
            
            # Handle forwarded emails (message/rfc822) - recursively extract their attachments
            # Do this BEFORE the is_attachment_part check since message/rfc822 parts should be processed specially
            if content_type == 'message/rfc822':
                try:
                    # Get the nested message
                    nested_payload = part.get_payload()
                    if nested_payload:
                        # nested_payload can be a list or a Message object
                        if isinstance(nested_payload, list):
                            for nested_msg in nested_payload:
                                if hasattr(nested_msg, 'walk'):
                                    logger.info(f"Processing nested message/rfc822 (depth={depth+1})")
                                    nested_texts, nested_infos = await self._extract_attachments(nested_msg, depth=depth+1)
                                    texts.extend(nested_texts)
                                    infos.extend(nested_infos)
                        elif hasattr(nested_payload, 'walk'):
                            logger.info(f"Processing nested message/rfc822 (depth={depth+1})")
                            nested_texts, nested_infos = await self._extract_attachments(nested_payload, depth=depth+1)
                            texts.extend(nested_texts)
                            infos.extend(nested_infos)
                except Exception as e:
                    logger.warning(f"Failed to process nested message/rfc822: {e}")
                continue  # Don't process the rfc822 part itself as an attachment
            
            # Use centralized attachment detection (handles forwarded emails, inline, document types)
            if is_attachment_part(part):
                # Handle unnamed attachments - generate filename from content-type
                # (filename and content_type already extracted above for detection)
                
                if not filename:
                    # Generate filename from content type
                    ext_map = {
                        'application/pdf': '.pdf',
                        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                        'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                        'text/plain': '.txt',
                        'text/calendar': '.ics',
                    }
                    ext = ext_map.get(content_type, '.bin')
                    filename = f"unnamed{ext}"
                    logger.debug(f"Generated filename for unnamed attachment: {filename}")
                
                # Clean filename - handle encoded filenames like .pdf?= or .pdf?utf-8?
                filename = self._clean_filename(filename)
                
                # Get payload
                try:
                    payload = part.get_payload(decode=True)
                except Exception as e:
                    logger.warning(f"Failed to decode attachment {filename}: {e}")
                    infos.append(AttachmentInfo(
                        filename=filename,
                        content_type=content_type,
                        size=0,
                        extracted_text=None,
                        extraction_error=f"Decode failed: {str(e)}"
                    ))
                    continue
                
                if not payload:
                    continue
                
                size = len(payload)
                
                # Check size limit
                if size > MAX_ATTACHMENT_SIZE:
                    logger.warning(f"Attachment {filename} ({size/1024/1024:.1f}MB) exceeds {MAX_ATTACHMENT_SIZE/1024/1024}MB limit, skipping extraction")
                    infos.append(AttachmentInfo(
                        filename=filename,
                        content_type=content_type,
                        size=size,
                        extracted_text=None,
                        extraction_error=f"File too large ({size/1024/1024:.1f}MB)"
                    ))
                    continue
                
                # Try to extract text based on type
                extracted_text = None
                error = None
                
                # Wrap each extraction in individual try-except for better isolation
                try:
                    if filename.lower().endswith('.pdf') or content_type == 'application/pdf':
                        try:
                            extracted_text = await self._extract_pdf_text(payload)
                        except Exception as e:
                            error = f"PDF: {str(e)[:100]}"
                    
                    elif filename.lower().endswith('.docx') or \
                         content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                        try:
                            extracted_text = await self._extract_docx_text(payload)
                        except Exception as e:
                            error = f"DOCX: {str(e)[:100]}"
                    
                    elif filename.lower().endswith(('.xlsx', '.xls')) or \
                         content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                        try:
                            extracted_text = await self._extract_excel_text(payload)
                        except Exception as e:
                            error = f"Excel: {str(e)[:100]}"
                    
                    elif filename.lower().endswith(('.pptx', '.ppt')) or \
                         content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                        try:
                            extracted_text = await self._extract_pptx_text(payload)
                        except Exception as e:
                            error = f"PPTX: {str(e)[:100]}"
                    
                    elif filename.lower().endswith('.rtf') or content_type == 'application/rtf':
                        try:
                            extracted_text = await self._extract_rtf_text(payload)
                        except Exception as e:
                            error = f"RTF: {str(e)[:100]}"
                    
                    elif filename.lower().endswith(('.ics', '.ical')) or content_type == 'text/calendar':
                        try:
                            extracted_text = await self._extract_ics_text(payload)
                        except Exception as e:
                            error = f"ICS: {str(e)[:100]}"
                    
                    elif filename.lower().endswith(('.txt', '.csv')):
                        try:
                            extracted_text = payload.decode('utf-8', errors='replace')
                        except Exception as e:
                            error = f"Text: {str(e)[:100]}"
                    
                    # MIME attachments are handled automatically by msg.walk()
                    # All attachment types are supported via content_disposition check
                    
                except Exception as e:
                    # Catch-all for unexpected errors during type determination
                    logger.error(f"Unexpected error processing {filename}: {e}")
                    error = f"Unexpected: {str(e)[:100]}"
                
                info = AttachmentInfo(
                    filename=filename,
                    content_type=content_type,
                    size=size,
                    extracted_text=extracted_text,
                    extraction_error=error
                )
                infos.append(info)
                
                # IMPORTANT: Always append to texts to keep alignment with infos
                # Empty string for failed extractions - prompt builder will skip these
                if extracted_text:
                    texts.append(f"[Attachment: {filename}]\n{extracted_text}")
                else:
                    texts.append("")  # Placeholder to maintain index alignment with infos
        
        return texts, infos
    
    async def _extract_pdf_text(self, pdf_bytes: bytes) -> str:
        """
        Extract text from PDF.
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_pdf(pdf_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("PDF", result.error or "unknown error", len(pdf_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        import warnings
        
        text_parts = []
        
        # Suppress pdfplumber warnings about invalid PDF structures
        # (common in PDFs with non-standard color definitions)
        with warnings.catch_warnings():
            warnings.filterwarnings('ignore', category=UserWarning)
            warnings.filterwarnings('ignore', message='.*invalid float value.*')
            
            try:
                with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                    for page in pdf.pages:
                        try:
                            text = page.extract_text()
                            if text:
                                text_parts.append(text)
                        except Exception as e:
                            logger.debug(f"Failed to extract text from PDF page: {e}")
                            continue
            except Exception as e:
                logger.warning(f"PDF extraction error: {e}")
                # Return empty rather than failing
                return ""
        
        return '\n\n'.join(text_parts)
    
    async def _extract_docx_text(self, docx_bytes: bytes) -> str:
        """
        Extract text from DOCX.
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_docx(docx_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("DOCX", result.error or "unknown error", len(docx_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        doc = Document(io.BytesIO(docx_bytes))
        text_parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
        return '\n'.join(text_parts)
    
    async def _extract_excel_text(self, xlsx_bytes: bytes) -> str:
        """
        Extract text from Excel files (XLSX).
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_xlsx(xlsx_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("XLSX", result.error or "unknown error", len(xlsx_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            text_parts.append(f"\n[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                # Skip empty rows
                if any(cell for cell in row):
                    row_text = '\t'.join(str(cell or '') for cell in row)
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    
    async def _extract_pptx_text(self, pptx_bytes: bytes) -> str:
        """
        Extract text from PowerPoint files (PPTX).
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_pptx(pptx_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("PPTX", result.error or "unknown error", len(pptx_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        prs = Presentation(io.BytesIO(pptx_bytes))
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n[Slide {i}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        
        return '\n'.join(text_parts)
    
    async def _extract_rtf_text(self, rtf_bytes: bytes) -> str:
        """
        Extract text from RTF files.
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_rtf(rtf_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("RTF", result.error or "unknown error", len(rtf_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        try:
            rtf_string = rtf_bytes.decode('utf-8', errors='replace')
            text = rtf_to_text(rtf_string)
            return text
        except Exception as e:
            logger.warning(f"Failed to parse RTF, trying latin-1 encoding: {e}")
            # Fallback to latin-1 encoding (common for older RTF files)
            rtf_string = rtf_bytes.decode('latin-1', errors='replace')
            return rtf_to_text(rtf_string)
    
    async def _extract_ics_text(self, ics_bytes: bytes) -> str:
        """
        Extract text from iCalendar (.ics) files.
        
        Security: When SANDBOX_ATTACHMENTS=true (default), runs in isolated subprocess.
        If sandbox fails, returns empty string (NO fallback to in-process for security).
        """
        if USE_SANDBOXED_EXTRACTION:
            result = await get_sandboxed_extractor().extract_ics(ics_bytes)
            if result.success:
                return result.text or ""
            else:
                _log_sandbox_failure("ICS", result.error or "unknown error", len(ics_bytes))
                return ""  # NO fallback - security over functionality
        
        # Direct extraction ONLY when sandboxing explicitly disabled
        try:
            ics_string = ics_bytes.decode('utf-8', errors='replace')
            cal = Calendar.from_ical(ics_string)
            
            text_parts = ["[Calendar Event]"]
            
            for component in cal.walk():
                if component.name == "VEVENT":
                    # Extract event details
                    summary = component.get('summary', 'No title')
                    start = component.get('dtstart')
                    end = component.get('dtend')
                    location = component.get('location', '')
                    description = component.get('description', '')
                    organizer = component.get('organizer', '')
                    
                    text_parts.append(f"Event: {summary}")
                    if start:
                        text_parts.append(f"Start: {start.dt}")
                    if end:
                        text_parts.append(f"End: {end.dt}")
                    if location:
                        text_parts.append(f"Location: {location}")
                    if organizer:
                        text_parts.append(f"Organizer: {organizer}")
                    if description:
                        text_parts.append(f"Description: {description}")
                    text_parts.append("")  # Blank line between events
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.warning(f"Failed to parse ICS file: {e}")
            # Fallback to raw text
            return ics_bytes.decode('utf-8', errors='replace')
    
    def _clean_filename(self, filename: str) -> str:
        """
        Clean filename - handle MIME encoded filenames.
        Handles cases like:
        - invoice.pdf?= 
        - document.pdf?utf-8?
        - file.docx?iso-8859-1?
        """
        if not filename:
            return filename
        
        # Remove MIME encoding artifacts (? characters and encodings)
        # Pattern: .ext?anything or .ext?=
        filename = re.sub(r'\?[^/]*$', '', filename)
        
        # Handle email.header.decode_header for encoded filenames
        from email.header import decode_header
        try:
            decoded = decode_header(filename)
            if decoded:
                # Take first part
                part, encoding = decoded[0]
                if isinstance(part, bytes):
                    filename = part.decode(encoding or 'utf-8', errors='replace')
                else:
                    filename = str(part)
        except:
            pass  # Keep original if decoding fails
        
        return filename.strip()

